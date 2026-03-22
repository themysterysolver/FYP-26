"""Build a single PPT slide with error_types_combined.csv as tables per dataset."""
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

HERE = Path(__file__).resolve().parent
CSV_PATH = HERE / "error_types_combined.csv"
OUT_PATH = HERE / "error_types_combined_FED-AVG-TT.pptx"


def main():
    df = pd.read_csv(CSV_PATH)
    datasets = df["dataset"].unique().tolist()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.45))
    p = title_box.text_frame.paragraphs[0]
    p.text = "FED-AVG-TT — Error types (baseline vs FedAvg)"
    p.font.size = Pt(24)
    p.font.bold = True

    col_w = Inches(4.0)
    gap = Inches(0.25)
    left0 = Inches(0.4)
    top = Inches(0.85)
    row_h = Inches(0.22)
    headers = ["Error type", "Baseline", "FedAvg", "Reduction"]

    for i, ds in enumerate(datasets):
        sub = df[df["dataset"] == ds].copy()
        nrows = len(sub) + 1
        left = left0 + i * (col_w + gap)

        table_shape = slide.shapes.add_table(
            nrows,
            len(headers),
            left,
            top,
            col_w,
            row_h * nrows,
        )
        tbl = table_shape.table

        for c, h in enumerate(headers):
            cell = tbl.cell(0, c)
            cell.text = h
            for para in cell.text_frame.paragraphs:
                para.font.bold = True
                para.font.size = Pt(10)

        for r, (_, row) in enumerate(sub.iterrows(), start=1):
            tbl.cell(r, 0).text = str(row["error_type"])
            tbl.cell(r, 1).text = str(int(row["baseline_count"]))
            tbl.cell(r, 2).text = str(int(row["fedavg_count"]))
            tbl.cell(r, 3).text = str(int(row["reduction"]))
            for c in range(4):
                for para in tbl.cell(r, c).text_frame.paragraphs:
                    para.font.size = Pt(9)

        cap = slide.shapes.add_textbox(left, top - Inches(0.28), col_w, Inches(0.25))
        cap.text_frame.paragraphs[0].text = str(ds)
        cap.text_frame.paragraphs[0].font.bold = True
        cap.text_frame.paragraphs[0].font.size = Pt(12)

    prs.save(OUT_PATH)
    print("Saved:", OUT_PATH)


if __name__ == "__main__":
    main()
